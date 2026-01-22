#!/usr/bin/env python3
"""
PX Center Branding Utilities
Brand constants and helper functions for creating PX-branded presentations,
documents, and spreadsheets.

Usage:
    from px_branding import create_branded_presentation, add_logo_to_slide
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# =============================================================================
# BRAND COLORS - PX CENTER OFFICIAL
# =============================================================================
# PX Center / Corporativo
PX_GRAY = RGBColor(0x6B, 0x77, 0x8C)  # #6B778C - Cinza PX Center

# Produtos
MOTORISTA_BLUE = RGBColor(0x29, 0x7A, 0xCC)   # #297ACC
AJUDANTE_GREEN = RGBColor(0x14, 0x8F, 0x4E)   # #148F4E
ACADEMIA_PURPLE = RGBColor(0x7D, 0x3D, 0xC5)  # #7D3DC5
RADAR_BLUE = RGBColor(0x1A, 0x4F, 0x80)       # #1A4F80
PULSAR_BLUE = RGBColor(0x42, 0x91, 0xFF)      # #4291FF
MASTERDRIVER_ORANGE = RGBColor(0xFA, 0x46, 0x16)  # #FA4616

# Neutros
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY_600 = RGBColor(0x4B, 0x55, 0x63)  # Body text

# Legado (para compatibilidade)
BLUE_100 = RGBColor(0xD5, 0xE4, 0xF5)  # Light backgrounds
BLUE_500 = RGBColor(0x33, 0x99, 0xFF)  # Primary, headers
BLUE_700 = RGBColor(0x0F, 0x3D, 0x66)  # Dark accents

# Color aliases
PRIMARY = PX_GRAY
PRIMARY_LIGHT = BLUE_100
PRIMARY_DARK = RADAR_BLUE

# =============================================================================
# TYPOGRAPHY
# =============================================================================
FONT_PRIMARY = 'Nunito'
FONT_SECONDARY = 'Poppins'
FONT_FALLBACK = 'Arial'

# Font sizes (in points)
FONT_SIZE_TITLE = 44
FONT_SIZE_SUBTITLE = 24
FONT_SIZE_HEADING = 32
FONT_SIZE_SUBHEADING = 24
FONT_SIZE_BODY = 18
FONT_SIZE_CAPTION = 14
FONT_SIZE_SMALL = 12

# =============================================================================
# SLIDE DIMENSIONS (16:9)
# =============================================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# =============================================================================
# PATHS (Windows compatible)
# =============================================================================
TEMP_DIR = os.path.join(os.environ.get('TEMP', '/tmp'), 'px-branding')

# Assets directory (relative to this script or skill root)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS_DIR = os.path.join(SKILL_DIR, 'assets')

# =============================================================================
# LOGO PATHS
# =============================================================================
LOGOS = {
    'px_preto': os.path.join(ASSETS_DIR, 'PX_preto (1).png'),
    'motorista_azul': os.path.join(ASSETS_DIR, 'Motorista-PX_azul (1).png'),
    'motorista_branco': os.path.join(ASSETS_DIR, 'Motorista-PX_branco.png'),
    'motorista_preto': os.path.join(ASSETS_DIR, 'Motorista-PX_preto.png'),
    'ajudante_verde': os.path.join(ASSETS_DIR, 'Ajudante-PX_verde_RGB (1).png'),
    'ajudante_branco': os.path.join(ASSETS_DIR, 'Ajudante-PX_branco.png'),
    'ajudante_preto': os.path.join(ASSETS_DIR, 'Ajudante-PX_preto.png'),
    'academia_roxo': os.path.join(ASSETS_DIR, 'Academia-PX_roxo_RGB (1).png'),
    'academia_branco': os.path.join(ASSETS_DIR, 'Academia-PX_branco.png'),
    'academia_preto': os.path.join(ASSETS_DIR, 'Academia-PX_preto.png'),
    'radar_azul': os.path.join(ASSETS_DIR, 'Radar-PX_azul (1).png'),
    'radar_branco': os.path.join(ASSETS_DIR, 'Radar-PX_branco.png'),
    'radar_preto': os.path.join(ASSETS_DIR, 'Radar-PX_preto.png'),
    'pulsar_azul': os.path.join(ASSETS_DIR, 'Pulsar_azul (1).png'),
    'pulsar_preto': os.path.join(ASSETS_DIR, 'Pulsar_preto_RGB.png'),
    'loja_px': os.path.join(ASSETS_DIR, 'LojaPX (1).png'),
}

def get_logo_path(logo_key='px_preto'):
    """Get path to a logo file. Returns None if not found."""
    path = LOGOS.get(logo_key)
    if path and os.path.exists(path):
        return path
    return None

# =============================================================================
# PRESENTATION HELPERS
# =============================================================================
def create_branded_presentation():
    """Create a new presentation with PX branding (16:9 aspect ratio)."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_logo_to_slide(slide, x=Inches(0.5), y=Inches(0.3), width=Inches(1.5), logo_key='px_preto'):
    """
    Add PX logo to a slide.

    Args:
        slide: The slide to add the logo to
        x: Left position (default: 0.5 inches from left)
        y: Top position (default: 0.3 inches from top)
        width: Logo width (default: 1.5 inches, height auto-calculated)
        logo_key: Which logo to use (default: 'px_preto')
                  Options: px_preto, motorista_azul/branco/preto,
                  ajudante_verde/branco/preto, academia_roxo/branco/preto,
                  radar_azul/branco/preto, pulsar_azul/preto, loja_px

    Returns:
        The picture shape or None if logo not found
    """
    logo_path = get_logo_path(logo_key)
    if logo_path:
        return slide.shapes.add_picture(logo_path, x, y, width=width)
    return None


def get_blank_slide_layout(prs):
    """Get the blank slide layout from a presentation."""
    return prs.slide_layouts[6]  # Typically blank


def add_title_slide(prs, title, subtitle=None):
    """
    Add a branded title slide.

    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Optional subtitle text

    Returns:
        The slide object
    """
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add background color
    background = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BLUE_700
    background.line.fill.background()

    # Add logo
    add_logo_to_slide(slide, x=Inches(0.5), y=Inches(0.3), width=Inches(2))

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(12.333), Inches(2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(FONT_SIZE_TITLE)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_PRIMARY

    # Add subtitle if provided
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5),
            Inches(12.333), Inches(1)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(FONT_SIZE_SUBTITLE)
        p.font.color.rgb = BLUE_100
        p.font.name = FONT_PRIMARY

    return slide


def add_content_slide(prs, title, add_logo=True):
    """
    Add a branded content slide with title.

    Args:
        prs: Presentation object
        title: Slide title
        add_logo: Whether to add logo (default: True)

    Returns:
        The slide object
    """
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add logo if requested
    if add_logo:
        add_logo_to_slide(slide, x=Inches(11.5), y=Inches(0.2), width=Inches(1.2))

    # Add title bar
    title_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLUE_700
    title_bar.line.fill.background()

    # Add title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(10), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(FONT_SIZE_HEADING)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_PRIMARY

    return slide


def style_text(paragraph, font_size=FONT_SIZE_BODY, bold=False, color=None, font_name=None):
    """
    Apply consistent styling to a paragraph.

    Args:
        paragraph: The paragraph to style
        font_size: Size in points (default: body size)
        bold: Whether to make bold (default: False)
        color: RGBColor (default: GRAY_600)
        font_name: Font name (default: FONT_PRIMARY)
    """
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color or GRAY_600
    paragraph.font.name = font_name or FONT_PRIMARY


# =============================================================================
# MAIN (for testing)
# =============================================================================
if __name__ == '__main__':
    print("PX Center Branding Utilities")
    print("=" * 40)
    print(f"Primary Color: #{BLUE_500}")
    print(f"Font: {FONT_PRIMARY}")
    print(f"Slide Size: {SLIDE_WIDTH} x {SLIDE_HEIGHT}")
    print()
    print("Creating test presentation...")

    prs = create_branded_presentation()
    add_title_slide(prs, "Test Presentation", "Created with px_branding.py")
    add_content_slide(prs, "Content Slide Example")

    output_path = os.path.join(TEMP_DIR, 'test_presentation.pptx')
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Saved to: {output_path}")
